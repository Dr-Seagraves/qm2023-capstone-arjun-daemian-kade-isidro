#!/usr/bin/env python3
"""Generate a PowerPoint presentation for the capstone project (CORRECTED VERSION)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_title_slide(prs, title, subtitle):
    """Add a title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(44, 62, 80)  # Dark blue/gray
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(236, 240, 241)
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(189, 195, 199)
    
    return slide

def add_content_slide(prs, title, content_points):
    """Add a content slide with bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(44, 62, 80)
    
    # Add a line
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(1.1), Inches(9.5), Inches(1.1))
    line.line.color.rgb = RGBColor(52, 152, 219)
    line.line.width = Pt(2)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.4), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(content_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        # Check if this is a sub-bullet (starts with "•")
        if point.startswith("  "):
            p.level = 1
            p.text = point.strip()
        else:
            p.level = 0
            p.text = point
        
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(44, 62, 80)
        p.space_before = Pt(4)
        p.space_after = Pt(4)
    
    return slide

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ==============================================================================
# SLIDE 1: Title
# ==============================================================================
add_title_slide(prs, 
    "Corruption, Risk & Foreign Investment",
    "QM 2023 Capstone Project\nArjun • Daemian • Kade • Isidro\nCorrected Analysis (2026-04-29)")

# ==============================================================================
# SLIDE 2: Refined Research Question
# ==============================================================================
add_content_slide(prs,
    "Research Question (Refined)",
    [
        "How do perceived corruption and geopolitical risk affect foreign direct investment?",
        "",
        "Originally Planned: Included crime data",
        "Actually Analyzed: CPI + GEPU (crime excluded due to time-invariance)",
        "",
        "Focus: Causal inference using panel econometric methods",
        "  with emphasis on methodological constraints",
    ])

# ==============================================================================
# SLIDE 3: Data Overview
# ==============================================================================
add_content_slide(prs,
    "Data Sources & Coverage",
    [
        "Sample: 20 countries, 13 years (2012-2024), 237 country-year observations",
        "",
        "Data Sources:",
        "  Corruption Perceptions Index (CPI) – Transparency International",
        "  Geopolitical Risk Index (GEPU) – Caldara & Iacoviello",
        "  Foreign Direct Investment (FDI) – World Bank",
        "  Crime Index – Numbeo (collected but see next slide)",
        "",
        "Outcome: FDI inflows (signed-log transformation)",
    ])

# ==============================================================================
# SLIDE 4: Crime Data Constraints
# ==============================================================================
add_content_slide(prs,
    "Crime Data: Why We Couldn't Use It",
    [
        "Problem: Crime is time-invariant within countries (e.g., US consistently ~40-45)",
        "",
        "Why This Matters in Fixed Effects Models:",
        "  Country fixed effects already capture all time-invariant differences",
        "  Crime became perfectly collinear with country dummies",
        "  Result: VIF = 6.56 (model fails to estimate)",
        "",
        "Solution: Excluded from main analysis",
        "  But kept in robustness checks and Random Forest",
        "  See CRIME_DATA_METHODOLOGY.md for full details",
    ])

# ==============================================================================
# SLIDE 5: Milestone 2 – Exploratory Data Analysis
# ==============================================================================
add_content_slide(prs,
    "M2: Exploratory Data Analysis",
    [
        "Approach: Correlation, time-series plots, group heterogeneity",
        "",
        "Key Findings:",
        "  Weak negative correlation: CPI vs FDI (r ≈ -0.14)",
        "  Suggested optimal lag: 12 years (for economic relationship)",
        "  China an outlier: much higher average CPI than peers",
        "  Crime 'not an important factor' in exploratory correlation",
        "",
        "Data Quality Flags:",
        "  Some countries missing 2025 data",
        "  China heterogeneity noted as potential issue",
    ])

# ==============================================================================
# SLIDE 6: Milestone 3 – Econometric Methodology
# ==============================================================================
add_content_slide(prs,
    "M3: Econometric Specification",
    [
        "Model: Two-way fixed effects (country + year dummies)",
        "",
        "Key Design Decisions:",
        "  Outcome: fdi_slog (signed-log transform; handles negative values)",
        "  Main driver: cpi_score_lag1 (one-year lag CPI)",
        "  Control: gepu (geopolitical uncertainty; time-varying)",
        "  Robust SEs: Clustered at country level",
        "",
        "Lag Specification Issue:",
        "  M2 identified lag=12 as optimal",
        "  But only 13 years of data ⟹ lag 12 = zero observations",
        "  Tested lags 1, 2, 3 as feasible alternatives",
    ])

# ==============================================================================
# SLIDE 7: M3 Main Results
# ==============================================================================
add_content_slide(prs,
    "M3: Primary Findings",
    [
        "Model A: Two-Way Fixed Effects",
        "",
        "CPI Effect (main hypothesis):",
        "  Coefficient: -0.1135  |  p-value: 0.4238  |  NOT SIGNIFICANT",
        "",
        "GEPU Effect (risk control):",
        "  Coefficient: -0.0159  |  p-value: 0.1876  |  NOT SIGNIFICANT",
        "",
        "Model Fit:",
        "  Within R² = 0.0149 (only 1.5% of variation explained)",
        "  Suggests major omitted drivers (macro conditions, trade policy, etc.)",
    ])

# ==============================================================================
# SLIDE 8: Critical Finding – 2020 Sensitivity
# ==============================================================================
add_content_slide(prs,
    "Critical Instability: 2020 Pandemic Shock",
    [
        "Baseline Model (with 2020): coef = -0.1135, p = 0.4238",
        "Excluding 2020 (just COVID year): coef = -0.0010, p = 0.9950",
        "",
        "What This Means:",
        "  99% coefficient reduction when pandemic year removed",
        "  Effect collapses to statistically zero",
        "  Suggests results dominated by COVID-era FDI shocks",
        "",
        "Implication:",
        "  Model captures crisis dynamics, NOT stable corruption-FDI relationship",
        "  Causal interpretation NOT justified",
        "  Need longer sample covering multiple normal economic regimes",
    ])

# ==============================================================================
# SLIDE 9: Robustness Checks
# ==============================================================================
add_content_slide(prs,
    "Robustness Checks: What Held Up?",
    [
        "Alternative Lags (1, 2, 3 years):",
        "  All negative, but all non-significant (p > 0.15)",
        "  Sign direction stable; magnitude not robust",
        "",
        "Crime-Based Subsamples:",
        "  Low-crime countries: coef = -0.168, p = 0.753",
        "  High-crime countries: coef = -0.001, p = 0.999",
        "  No significant heterogeneity detected",
        "",
        "Clustered vs. Unclustered SEs:",
        "  Conclusion unchanged under both covariance assumptions",
    ])

# ==============================================================================
# SLIDE 10: Model Limitations
# ==============================================================================
add_content_slide(prs,
    "Limitations & Caveats",
    [
        "Omitted Variables (major FDI drivers missing):",
        "  GDP growth, exchange rates, interest rates, commodity prices",
        "  Trade openness, infrastructure quality, tax policy",
        "",
        "Measurement Error:",
        "  CPI perception-based; Crime data Numbeo-sourced; GEPU broad",
        "",
        "External Validity:",
        "  Only 20 countries with complete overlap",
        "  Results may not generalize to other regions/periods",
        "",
        "Endogeneity:",
        "  FDI may improve governance, not just vice versa",
    ])

# ==============================================================================
# SLIDE 11: Conclusions & Main Takeaways
# ==============================================================================
add_content_slide(prs,
    "Conclusions",
    [
        "Finding: Corruption (CPI) effect on FDI is NOT statistically significant",
        "  and highly unstable (driven by 2020 pandemic shock)",
        "",
        "Policy Implication:",
        "  Corruption reform matters for governance, but:",
        "  FDI also depends on macro stability, uncertainty, and broader institutions",
        "  Anti-corruption alone won't substantially increase investment",
        "",
        "Methodological Learning:",
        "  Criminal data collected but inapplicable in FE framework (time-invariant)",
        "  Data constraints forced lag compromise (wanted 12, got 1)",
        "  Short panel with major shocks = weak causal inference",
    ])

# ==============================================================================
# SLIDE 12: Future Research Directions
# ==============================================================================
add_content_slide(prs,
    "Recommendations for Future Work",
    [
        "Extend Data:",
        "  Add 15-20 years of historical data (pre-2012) to test lag-12 relationship",
        "  Reduce relative weight of 2020 crisis year",
        "",
        "Add Macro Controls:",
        "  Include real interest rates, exchange rates, GDP growth, commodity prices",
        "",
        "Alternative Methods:",
        "  Use instrumental variables if valid instruments found",
        "  Run Random Effects (RE) models to estimate crime effects",
        "  Investigate 2020 structural break separately",
        "",
        "Qualitative Research:",
        "  Interview investors on actual corruption-FDI linkage",
    ])

# ==============================================================================
# SLIDE 13: Takeaway Message
# ==============================================================================
add_content_slide(prs,
    "Key Takeaway",
    [
        "What We Learned:",
        "  Corruption matters qualitatively for investor confidence",
        "  But quantitative causal effect NOT detected in this short panel",
        "",
        "Why?",
        "  Pandemic dominated recent decade",
        "  Missing macro drivers (interest rates, growth, trade)",
        "  Optimal lag (12 years) unattainable with 13-year sample",
        "",
        "Next Steps:",
        "  Get longer data + macro controls",
        "  Then revisit causal question with more confidence",
        "  This analysis is exploratory foundation, not final answer",
    ])

# ==============================================================================
# Save the presentation
# ==============================================================================
prs.save("capstone_presentation_corrected.pptx")
print("✓ Presentation saved: capstone_presentation_corrected.pptx (13 slides)")
